# 模块用途：解析 .docx、.pdf、.pptx 文件，提取纯文本内容
# .docx/.pptx 额外支持：
#   1. 调用本机 Office（win32com）把旧式 OLE 公式转换为现代 OMML 格式
#   2. 把 OMML 数学公式转换为 $...$ / $$...$$ LaTeX 格式

import os
import tempfile
import shutil

# ── OMML 命名空间 ──────────────────────────────
_M = 'http://schemas.openxmlformats.org/officeDocument/2006/math'
_W = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'


def _tag(elem) -> str:
    return elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag


def _ns(elem) -> str:
    return elem.tag.split('}')[0][1:] if '}' in elem.tag else ''


def _mval(elem) -> str:
    return elem.get(f'{{{_M}}}val') or elem.get('val') or ''


# ──────────────────────────────────────────────
# win32com：调用本机 Office 转换旧式 OLE 公式
# ──────────────────────────────────────────────

def _convert_with_office(filepath: str, filetype: str) -> str:
    """
    用本机 Word / PowerPoint 打开文件，转换为现代格式后存为临时文件。
    返回临时文件路径（调用方负责删除）。
    转换失败（未安装 Office 等）时返回原路径，不抛异常。

    filetype: 'word' | 'pptx'
    """
    try:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()

        # 必须用绝对路径，否则 COM 找不到文件
        abs_path = os.path.abspath(filepath)
        ext = os.path.splitext(filepath)[1].lower()
        tmp_path = os.path.join(
            tempfile.gettempdir(),
            f'_crow_converted_{os.getpid()}{ext}'
        )

        if filetype == 'word':
            word = win32com.client.Dispatch('Word.Application')
            word.Visible = False
            word.DisplayAlerts = False
            try:
                doc = word.Documents.Open(abs_path, ReadOnly=True)
                # Convert() 把整个文档（含 OLE 公式）升级为现代格式
                doc.Convert()
                # 16 = wdFormatDocumentDefault (.docx)
                doc.SaveAs2(tmp_path, FileFormat=16)
                doc.Close(SaveChanges=False)
            finally:
                word.Quit()

        else:  # pptx
            ppt = win32com.client.Dispatch('PowerPoint.Application')
            ppt.Visible = False
            try:
                prs = ppt.Presentations.Open(abs_path, ReadOnly=True,
                                             Untitled=False, WithWindow=False)
                # 24 = ppSaveAsOpenXMLPresentation (.pptx)
                prs.SaveAs(tmp_path, 24)
                prs.Close()
            finally:
                ppt.Quit()

        pythoncom.CoUninitialize()

        if os.path.exists(tmp_path):
            return tmp_path
        return filepath

    except Exception:
        # 未安装 Office / COM 失败 → 降级，用原文件继续解析
        return filepath


# ──────────────────────────────────────────────
# OMML → LaTeX 转换器
# ──────────────────────────────────────────────

def _omml_to_latex(omath_elem) -> str:
    _PROP_TAGS = {
        'oMathPr', 'rPr', 'fPr', 'radPr', 'sSupPr', 'sSubPr', 'sSubSupPr',
        'naryPr', 'dPr', 'funcPr', 'mPr', 'mrPr', 'limLowPr', 'limUppPr',
        'eqArrPr', 'groupChrPr', 'phantPr', 'phant', 'ctrlPr',
    }
    _FUNC_MAP = {
        'sin': r'\sin', 'cos': r'\cos', 'tan': r'\tan', 'cot': r'\cot',
        'sec': r'\sec', 'csc': r'\csc', 'arcsin': r'\arcsin',
        'arccos': r'\arccos', 'arctan': r'\arctan',
        'ln': r'\ln', 'log': r'\log', 'exp': r'\exp',
        'lim': r'\lim', 'max': r'\max', 'min': r'\min',
        'sup': r'\sup', 'inf': r'\inf', 'det': r'\det',
    }
    _NARY_MAP = {
        '∫': r'\int', '∬': r'\iint', '∭': r'\iiint',
        '∮': r'\oint', '∑': r'\sum', '∏': r'\prod',
        '⋃': r'\bigcup', '⋂': r'\bigcap',
    }

    def children_latex(elem) -> str:
        return ''.join(convert(c) for c in elem)

    def child_by_tag(elem, local) -> str:
        for c in elem:
            if _tag(c) == local:
                return children_latex(c)
        return ''

    def convert(elem) -> str:
        t = _tag(elem)

        if t in _PROP_TAGS:
            return ''
        if t in ('oMath', 'oMathPara', 'e', 'num', 'den',
                 'deg', 'lim', 'sup', 'sub', 'fName'):
            return children_latex(elem)
        if t == 'r':
            return ''.join(c.text or '' for c in elem if _tag(c) == 't')
        if t == 'f':
            return rf'\frac{{{child_by_tag(elem, "num")}}}{{{child_by_tag(elem, "den")}}}'
        if t == 'rad':
            deg = child_by_tag(elem, 'deg').strip()
            e   = child_by_tag(elem, 'e')
            return rf'\sqrt[{deg}]{{{e}}}' if deg else rf'\sqrt{{{e}}}'
        if t == 'sSup':
            return rf'{child_by_tag(elem, "e")}^{{{child_by_tag(elem, "sup")}}}'
        if t == 'sSub':
            return rf'{child_by_tag(elem, "e")}_{{{child_by_tag(elem, "sub")}}}'
        if t == 'sSubSup':
            return (rf'{child_by_tag(elem, "e")}'
                    rf'_{{{child_by_tag(elem, "sub")}}}'
                    rf'^{{{child_by_tag(elem, "sup")}}}')
        if t == 'nary':
            sym = r'\int'
            sub = sup = e = ''
            for c in elem:
                ct = _tag(c)
                if ct == 'naryPr':
                    for pr in c:
                        if _tag(pr) == 'chr':
                            sym = _NARY_MAP.get(_mval(pr), r'\int')
                elif ct == 'sub':
                    sub = children_latex(c)
                elif ct == 'sup':
                    sup = children_latex(c)
                elif ct == 'e':
                    e = children_latex(c)
            result = sym
            if sub: result += rf'_{{{sub}}}'
            if sup: result += rf'^{{{sup}}}'
            return result + f' {e}'
        if t == 'd':
            open_ch = '('; close_ch = ')'
            for c in elem:
                if _tag(c) == 'dPr':
                    for pr in c:
                        pt = _tag(pr)
                        if pt == 'begChr': open_ch = _mval(pr) or '('
                        elif pt == 'endChr': close_ch = _mval(pr) or ')'
            _lmap = {'(':'(','[':'[','{':r'\{','|':r'|','':r'.'}
            _rmap = {')':')',']':']','}':r'\}','|':r'|','':r'.'}
            contents = [children_latex(c) for c in elem if _tag(c) == 'e']
            return (rf'\left{_lmap.get(open_ch, open_ch)}'
                    + ', '.join(contents)
                    + rf'\right{_rmap.get(close_ch, close_ch)}')
        if t == 'func':
            fname = child_by_tag(elem, 'fName').strip()
            return rf'{_FUNC_MAP.get(fname, fname)} {{{child_by_tag(elem, "e")}}}'
        if t == 'limLow':
            return rf'{child_by_tag(elem, "e")}_{{{child_by_tag(elem, "lim")}}}'
        if t == 'limUpp':
            return rf'{child_by_tag(elem, "e")}^{{{child_by_tag(elem, "lim")}}}'
        if t == 'm':
            rows = []
            for c in elem:
                if _tag(c) == 'mr':
                    rows.append(' & '.join(children_latex(cell)
                                           for cell in c if _tag(cell) == 'e'))
            return r'\begin{matrix}' + r' \\ '.join(rows) + r'\end{matrix}'
        if t == 'groupChr':
            chr_val = ''; e = ''
            for c in elem:
                ct = _tag(c)
                if ct == 'groupChrPr':
                    for pr in c:
                        if _tag(pr) == 'chr': chr_val = _mval(pr)
                elif ct == 'e':
                    e = children_latex(c)
            _gcmap = {'̄': r'\overline', '̂': r'\hat', '̃': r'\tilde',
                      '⏞': r'\overbrace', '⏟': r'\underbrace'}
            return rf'{_gcmap.get(chr_val, r"\overline")}{{{e}}}'
        if t == 'eqArr':
            rows = [children_latex(c) for c in elem if _tag(c) == 'e']
            return r'\begin{aligned}' + r' \\ '.join(rows) + r'\end{aligned}'
        return children_latex(elem)

    return convert(omath_elem)


# ──────────────────────────────────────────────
# 段落文本提取（含公式）
# ──────────────────────────────────────────────

def _extract_paragraph_text(para) -> str:
    has_math = False
    parts = []

    def process(elem):
        nonlocal has_math
        t  = _tag(elem)
        ns = _ns(elem)

        if t == 'oMathPara' and ns == _M:
            has_math = True
            for c in elem:
                if _tag(c) == 'oMath':
                    latex = _omml_to_latex(c).strip()
                    parts.append(f'$${latex}$$')

        elif t == 'oMath' and ns == _M:
            has_math = True
            latex = _omml_to_latex(elem).strip()
            parts.append(f'${latex}$')

        elif t == 'object':
            # OLE 嵌入对象（旧式 Equation Editor 3.0）
            # 若 win32com 转换成功，这里不会再出现；
            # 转换失败时降级显示占位符
            _O = 'urn:schemas-microsoft-com:office:office'
            is_equation = any(
                'Equation' in (c.get(f'{{{_O}}}ProgID') or c.get('ProgID') or '')
                for c in elem
            )
            if is_equation:
                has_math = True
                parts.append('[公式]')

        elif t == 'r' and ns == _W:
            has_ole = any(_tag(c) == 'object' for c in elem)
            if not has_ole:
                for c in elem:
                    if _tag(c) == 't':
                        parts.append(c.text or '')

        elif t == 'pPr':
            pass

        else:
            for c in elem:
                process(c)

    for child in para._p:
        process(child)

    if has_math:
        return ''.join(parts)
    return para.text or ''


# ──────────────────────────────────────────────
# 公开解析接口
# ──────────────────────────────────────────────

def parse_docx(filepath: str) -> str:
    """解析 Word 文档，返回纯文本（公式转为 LaTeX）。
    自动尝试用本机 Word 转换旧式 OLE 公式，失败时降级为 [公式] 占位符。"""
    tmp = None
    try:
        from docx import Document

        # 尝试用 Word 转换（处理旧式 OLE 公式）
        converted = _convert_with_office(filepath, 'word')
        if converted != filepath:
            tmp = converted  # 记录临时文件，最后清理

        doc = Document(converted)
        paragraphs = []
        for p in doc.paragraphs:
            text = _extract_paragraph_text(p).strip()
            if text:
                paragraphs.append(text)
        return '\n'.join(paragraphs)

    except Exception as e:
        raise ValueError(f'无法读取 Word 文档：{e}')
    finally:
        if tmp and os.path.exists(tmp):
            try:
                os.remove(tmp)
            except Exception:
                pass


def parse_pdf(filepath: str) -> str:
    """解析 PDF 文件，返回纯文本"""
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    texts.append(text.strip())
        if not texts:
            raise ValueError('PDF 中未能提取到文字，可能是扫描件图片格式')
        return '\n'.join(texts)
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f'无法读取 PDF 文件：{e}')


def parse_pptx(filepath: str) -> str:
    """解析 PPT 文件，每张幻灯片标题作为小节名，正文文本拼合为内容。
    自动尝试用本机 PowerPoint 转换旧式 OLE 公式，失败时降级为 [公式] 占位符。"""
    tmp = None
    try:
        from pptx import Presentation

        # 尝试用 PowerPoint 转换
        converted = _convert_with_office(filepath, 'pptx')
        if converted != filepath:
            tmp = converted

        prs = Presentation(converted)
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            title = ''
            bodies = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape.shape_type == 13 or 'title' in shape.name.lower():
                    title = shape.text_frame.text.strip()
                else:
                    text = shape.text_frame.text.strip()
                    if text:
                        bodies.append(text)
            if not title:
                title = f'第 {i} 页'
            block = title + '\n' + '\n'.join(bodies) if bodies else title
            slides_text.append(block)
        if not slides_text:
            raise ValueError('PPT 中未能提取到文字内容')
        return '\n\n'.join(slides_text)

    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f'无法读取 PPT 文件：{e}')
    finally:
        if tmp and os.path.exists(tmp):
            try:
                os.remove(tmp)
            except Exception:
                pass


def parse_file(filepath: str) -> str:
    """根据文件扩展名自动选择解析方式，返回纯文本"""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.docx':
        return parse_docx(filepath)
    elif ext == '.pdf':
        return parse_pdf(filepath)
    elif ext == '.pptx':
        return parse_pptx(filepath)
    else:
        raise ValueError(f'不支持的文件格式：{ext}，请上传 .docx、.pdf 或 .pptx 文件')
